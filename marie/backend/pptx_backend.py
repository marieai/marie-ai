import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from marie.backend.base_backend import DocumentBackend

_log = logging.getLogger(__name__)

# EMU (English Metric Units) to points: 1 pt = 12700 EMU
_EMU_TO_PT = 1.0 / 12700.0


def _emu_to_pt(emu) -> float:
    if emu is None:
        return 0.0
    return float(emu) * _EMU_TO_PT


def _text_to_words(text: str, x0: float, y0: float, x1: float, y1: float) -> list[dict]:
    parts = text.split()
    if not parts:
        return []
    total_chars = sum(len(w) for w in parts)
    if total_chars == 0:
        return []
    span = x1 - x0
    words = []
    cx = x0
    for w in parts:
        frac = len(w) / total_chars
        wx1 = cx + span * frac
        words.append(
            {
                "text": w,
                "bbox": [round(cx, 2), round(y0, 2), round(wx1, 2), round(y1, 2)],
                "confidence": 1.0,
            }
        )
        cx = wx1
    return words


def _extract_shape_text(shape, slide_w: float, slide_h: float) -> list[dict]:
    """Extract lines from a single shape, handling text frames and tables."""
    lines: list[dict] = []

    left = _emu_to_pt(shape.left) if shape.left is not None else 0.0
    top = _emu_to_pt(shape.top) if shape.top is not None else 0.0
    width = _emu_to_pt(shape.width) if shape.width is not None else slide_w
    height = _emu_to_pt(shape.height) if shape.height is not None else slide_h

    # Tables
    if shape.has_table:
        table = shape.table
        row_h = height / max(len(table.rows), 1)
        for row_idx, row in enumerate(table.rows):
            col_w = width / max(len(row.cells), 1)
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip()
                if not cell_text:
                    continue
                cx0 = left + col_idx * col_w
                cy0 = top + row_idx * row_h
                cx1 = cx0 + col_w
                cy1 = cy0 + row_h
                words = _text_to_words(cell_text, cx0, cy0, cx1, cy1)
                if words:
                    lines.append(
                        {
                            "text": cell_text,
                            "bbox": [
                                round(cx0, 2),
                                round(cy0, 2),
                                round(cx1, 2),
                                round(cy1, 2),
                            ],
                            "words": words,
                        }
                    )
        return lines

    # Text frames
    if not hasattr(shape, "text_frame"):
        return lines

    try:
        tf = shape.text_frame
    except Exception:
        return lines

    n_paras = len(tf.paragraphs)
    para_h = height / max(n_paras, 1)

    for i, paragraph in enumerate(tf.paragraphs):
        text = paragraph.text.strip()
        if not text:
            continue
        y0 = top + i * para_h
        y1 = y0 + para_h
        x0, x1 = left, left + width
        words = _text_to_words(text, x0, y0, x1, y1)
        if words:
            lines.append(
                {
                    "text": text,
                    "bbox": [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)],
                    "words": words,
                }
            )

    return lines


def _walk_shapes(shapes, slide_w: float, slide_h: float) -> list[dict]:
    """Recursively walk shapes including grouped shapes."""
    lines: list[dict] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            lines.extend(_walk_shapes(shape.shapes, slide_w, slide_h))
        else:
            lines.extend(_extract_shape_text(shape, slide_w, slide_h))
    return lines


class PptxBackend(DocumentBackend):
    @classmethod
    def supported_formats(cls) -> set[str]:
        return {"pptx"}

    @classmethod
    def is_available(cls) -> bool:
        try:
            import pptx  # noqa: F401

            return True
        except ImportError:
            return False

    def convert(self, file_path: str, **kwargs) -> dict:
        prs = Presentation(file_path)

        slide_w = _emu_to_pt(prs.slide_width)
        slide_h = _emu_to_pt(prs.slide_height)

        results: list[dict] = []
        for slide_idx, slide in enumerate(prs.slides):
            lines = _walk_shapes(slide.shapes, slide_w, slide_h)
            result = {
                "words": [w for line in lines for w in line["words"]],
                "lines": lines,
                "meta": {"page": slide_idx, "width": slide_w, "height": slide_h},
            }
            results.append(result)

        return {"mode": "parsed", "results": results, "pages": len(results)}
